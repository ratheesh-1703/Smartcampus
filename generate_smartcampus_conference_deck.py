from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "SmartCampus_Conference_Deck.pptx"
MAP_IMAGE = ROOT / "frontend" / "public" / "kalasalingam-map.png"

NAVY = RGBColor(11, 22, 42)
MIDNIGHT = RGBColor(17, 34, 61)
GOLD = RGBColor(214, 173, 74)
TEAL = RGBColor(58, 188, 194)
SAND = RGBColor(245, 241, 232)
MIST = RGBColor(236, 241, 247)
SLATE = RGBColor(91, 103, 124)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(20, 24, 31)
RED = RGBColor(194, 72, 72)
GREEN = RGBColor(58, 139, 92)


def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def style_paragraph(paragraph, font_size=20, color=BLACK, bold=False, bullet=False, level=0, name="Aptos"):
    paragraph.font.name = name
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.level = level
    paragraph.bullet = bullet


def add_rounded_card(slide, left, top, width, height, fill, line=None):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line or fill
    return card


def add_textbox(slide, left, top, width, height, text, font_size=20, color=BLACK, bold=False, align=PP_ALIGN.LEFT, fill=None, font_name="Aptos"):
    shape = slide.shapes.add_textbox(left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shape


def add_bullet_box(slide, left, top, width, height, title, bullets, title_color=NAVY, body_color=BLACK, fill=MIST):
    box = add_rounded_card(slide, left, top, width, height, fill)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = title_color

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        style_paragraph(p, font_size=14, color=body_color, bullet=True)
    return box


def add_footer(slide, text):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.333), Inches(0.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.color.rgb = GOLD
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(7.06), Inches(12.3), Inches(0.25))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = SLATE


def add_header(slide, title, subtitle=None):
    ribbon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.88))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = NAVY
    ribbon.line.color.rgb = NAVY
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.88), Inches(13.333), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.color.rgb = GOLD

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(8.4), Inches(0.42))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.47), Inches(0.49), Inches(8.6), Inches(0.2))
        tf = sub.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(10.5)
        r.font.color.rgb = RGBColor(202, 211, 225)


def add_tag(slide, text, left, top, width=1.35, color=GOLD, text_color=NAVY):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(width), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = text_color
    return shape


def add_connection(slide, x1, y1, x2, y2, color=SLATE, width=1.75):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NAVY)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    band.fill.solid()
    band.fill.fore_color.rgb = MIDNIGHT
    band.line.color.rgb = MIDNIGHT

    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    side.fill.solid()
    side.fill.fore_color.rgb = GOLD
    side.line.color.rgb = GOLD

    title = slide.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(6.8), Inches(1.8))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "SmartCampus"
    r.font.name = "Aptos Display"
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "Role-aware campus operations platform"
    style_paragraph(p2, font_size=20, color=RGBColor(228, 234, 243))
    p3 = tf.add_paragraph()
    p3.text = "Attendance integrity, emergency response, live location intelligence, and institutional workflows"
    style_paragraph(p3, font_size=16, color=RGBColor(206, 216, 230))

    add_tag(slide, "Attendance Integrity", Inches(0.7), Inches(3.35), 1.75, color=GOLD)
    add_tag(slide, "Campus Safety", Inches(2.55), Inches(3.35), 1.45, color=TEAL, text_color=WHITE)
    add_tag(slide, "Role-Based Ops", Inches(4.1), Inches(3.35), 1.6, color=SAND, text_color=NAVY)

    card = add_rounded_card(slide, Inches(7.75), Inches(1.1), Inches(4.8), Inches(5.4), WHITE)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Conference focus"
    r.font.name = "Aptos Display"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = NAVY
    for line in [
        "Unified platform replacing siloed attendance, safety, and office tools",
        "PHP API layer with token-based access and role-scoped data retrieval",
        "React frontend with role dashboards, Leaflet maps, and near-real-time polling",
        "Attendance workflow includes QR + hotspot/network-prefix verification",
        "SOS and live-location flows are wired for escalation and governance",
    ]:
        p = tf.add_paragraph()
        p.text = line
        style_paragraph(p, font_size=14, color=BLACK, bullet=True)

    add_footer(slide, "SmartCampus conference deck | generated from repository analysis")


def problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MIST)
    add_header(slide, "Why SmartCampus Exists", "The platform addresses fragmented campus operations and weak cross-module accountability")

    cards = [
        ("Data silos", "Separate tools for attendance, finance, safety, and records duplicate effort and fragment truth.", RED),
        ("Operational delay", "Manual reconciliation slows attendance closure, notifications, and incident handling.", GOLD),
        ("Safety gaps", "No single workflow for SOS, live location, and escalation across roles.", TEAL),
        ("Inconsistent access", "Responsibilities differ by role, but many systems do not enforce that boundary cleanly.", NAVY),
    ]
    x = 0.55
    y = 1.25
    for i, (title, body, color) in enumerate(cards):
        card = add_rounded_card(slide, Inches(x + (i % 2) * 6.25), Inches(y + (i // 2) * 1.65), Inches(5.55), Inches(1.35), WHITE)
        accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + (i % 2) * 6.25), Inches(y + (i // 2) * 1.65), Inches(0.15), Inches(1.35))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.color.rgb = color
        tf = card.text_frame
        tf.clear()
        tf.margin_left = Pt(16)
        tf.margin_right = Pt(12)
        tf.margin_top = Pt(10)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = "Aptos Display"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.text = body
        style_paragraph(p2, font_size=12.5, color=BLACK)

    add_bullet_box(slide, Inches(0.65), Inches(4.85), Inches(6.0), Inches(1.5), "Core response", [
        "One frontend, one backend, one data model",
        "Role-aware routing and backend authorization",
        "Integrity controls where the workflow is security-sensitive",
    ])
    add_bullet_box(slide, Inches(6.85), Inches(4.85), Inches(5.8), Inches(1.5), "Result", [
        "Less duplication",
        "Faster handoff between academic and admin teams",
        "Clear audit trail for attendance and safety actions",
    ])
    add_footer(slide, "SmartCampus conference deck | problem statement")


def architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_header(slide, "System Architecture", "React frontend, PHP service layer, and MySQL data persistence")

    layers = [
        ("React frontend", "Role-based dashboards\nVite + React Router\nBootstrap, Leaflet, Chart.js", TEAL),
        ("PHP API layer", "Modular endpoints\nAuth, validation, CORS, security headers\nBearer token checks", NAVY),
        ("MySQL backend", "Users, sessions, attendance, SOS, live locations\nInstitutional module tables", GOLD),
    ]
    for idx, (title, body, color) in enumerate(layers):
        add_rounded_card(slide, Inches(0.9), Inches(1.4 + idx * 1.45), Inches(5.0), Inches(1.05), MIST)
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.49 + idx * 1.45), Inches(1.6), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.color.rgb = color
        tf = badge.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        r.font.name = "Aptos"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = WHITE if color != GOLD else NAVY

        t = slide.shapes.add_textbox(Inches(2.72), Inches(1.42 + idx * 1.45), Inches(3.0), Inches(0.78))
        tf = t.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = body
        style_paragraph(p, font_size=12.5, color=BLACK)

    add_connection(slide, Inches(6.1), Inches(1.95), Inches(7.15), Inches(1.95), color=SLATE, width=2.2)
    add_connection(slide, Inches(6.1), Inches(3.35), Inches(7.15), Inches(3.35), color=SLATE, width=2.2)
    add_connection(slide, Inches(6.1), Inches(4.8), Inches(7.15), Inches(4.8), color=SLATE, width=2.2)

    right = add_rounded_card(slide, Inches(7.3), Inches(1.25), Inches(5.5), Inches(4.55), NAVY)
    tf = right.text_frame
    tf.clear()
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Request lifecycle"
    r.font.name = "Aptos Display"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WHITE
    for line in [
        "1. User logs in and receives a bearer token",
        "2. Frontend attaches Authorization header automatically",
        "3. Backend validates token, role, and scope",
        "4. Endpoint executes domain logic and returns JSON",
        "5. Role-specific dashboard renders the result",
    ]:
        p = tf.add_paragraph()
        p.text = line
        style_paragraph(p, font_size=13, color=RGBColor(236, 240, 247), bullet=True)

    add_bullet_box(slide, Inches(7.3), Inches(5.95), Inches(5.5), Inches(0.7), "Evidence", [
        "React/Vite frontend, PHP backend, MySQL schema, and bearer-token auth are documented across the repo.",
    ], title_color=NAVY, body_color=BLACK, fill=SAND)
    add_footer(slide, "SmartCampus conference deck | architecture")


def modules_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MIST)
    add_header(slide, "Modules and Roles", "A single platform spans academic, operational, and service functions")

    left = add_rounded_card(slide, Inches(0.6), Inches(1.25), Inches(5.8), Inches(5.35), WHITE)
    tf = left.text_frame
    tf.clear()
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Major modules"
    r.font.name = "Aptos Display"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = NAVY
    for item in [
        "Academics: student/teacher management, timetables, marks, subject assignment",
        "Attendance: session creation, marking, validation, rejection logs",
        "Safety: SOS reporting, live location, escalation intelligence",
        "Registrar: records, certificates, admissions workflows",
        "Finance: fees, payments, budgets, cost centers",
        "Exam, placement, hostel, library, dean, affairs, parent services",
    ]:
        p = tf.add_paragraph()
        p.text = item
        style_paragraph(p, font_size=13.5, color=BLACK, bullet=True)

    right = add_rounded_card(slide, Inches(6.6), Inches(1.25), Inches(6.0), Inches(5.35), NAVY)
    tf = right.text_frame
    tf.clear()
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Role coverage"
    r.font.name = "Aptos Display"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WHITE
    for item in [
        "Admin and affairs: global oversight and incident visibility",
        "Teacher, HOD, coordinator: scoped academic and attendance actions",
        "Student and parent: self-service access and notifications",
        "Accountant, registrar, exam controller: office workflows",
        "Placement, hostel warden, librarian, dean: role-specific dashboards",
    ]:
        p = tf.add_paragraph()
        p.text = item
        style_paragraph(p, font_size=13.5, color=RGBColor(236, 240, 247), bullet=True)

    highlight = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.1), Inches(6.0), Inches(9.1), Inches(0.6))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = SAND
    highlight.line.color.rgb = SAND
    tf = highlight.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "The platform is designed around role boundaries, not just features."
    r.font.name = "Aptos Display"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = NAVY

    add_footer(slide, "SmartCampus conference deck | modules and roles")


def attendance_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_header(slide, "Attendance Integrity", "A multi-layer verification pipeline blocks proxy attendance and records rejections")

    x_positions = [0.55, 2.45, 4.35, 6.25, 8.15]
    labels = [
        ("Teacher starts", "Session context, class, period, network prefix, QR token"),
        ("Student joins", "Available sessions are filtered by class membership"),
        ("Token check", "Bearer token, active session, and student role are validated"),
        ("Network gate", "Gateway prefix comparison detects hotspot mismatch"),
        ("Write audit", "Success inserts attendance; failure writes rejection reason"),
    ]
    colors = [TEAL, NAVY, GOLD, RED, GREEN]
    for idx, ((title, body), x, color) in enumerate(zip(labels, x_positions, colors)):
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(2.0), Inches(1.6), Inches(1.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.color.rgb = color
        tf = circle.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(idx + 1)
        r.font.name = "Aptos Display"
        r.font.size = Pt(26)
        r.font.bold = True
        r.font.color.rgb = WHITE

        t = slide.shapes.add_textbox(Inches(x - 0.25), Inches(3.8), Inches(2.1), Inches(1.1))
        tf = t.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        r.font.name = "Aptos Display"
        r.font.size = Pt(14.5)
        r.font.bold = True
        r.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.text = body
        style_paragraph(p2, font_size=11.5, color=BLACK)
        p2.alignment = PP_ALIGN.CENTER

    for start_x, end_x in zip(x_positions[:-1], x_positions[1:]):
        add_connection(slide, Inches(start_x + 1.6), Inches(2.8), Inches(end_x), Inches(2.8), color=SLATE, width=2.0)

    add_bullet_box(slide, Inches(0.7), Inches(5.2), Inches(5.8), Inches(1.4), "Verification layers", [
        "QR token match",
        "Network prefix check",
        "Duplicate attendance prevention",
        "Rejection log for suspicious attempts",
    ], title_color=NAVY, body_color=BLACK, fill=MIST)
    add_bullet_box(slide, Inches(6.8), Inches(5.2), Inches(5.9), Inches(1.4), "Conference takeaway", [
        "Attendance is not just recorded; it is verified, scoped, and audited.",
    ], title_color=NAVY, body_color=BLACK, fill=SAND)
    add_footer(slide, "SmartCampus conference deck | attendance integrity")


def safety_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MIST)
    add_header(slide, "Safety, SOS, and Live Location", "Emergency reporting is paired with real-time movement monitoring and scoped visibility")

    left = add_rounded_card(slide, Inches(0.65), Inches(1.25), Inches(5.4), Inches(5.25), WHITE)
    tf = left.text_frame
    tf.clear()
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "SOS workflow"
    r.font.name = "Aptos Display"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = NAVY
    for item in [
        "Student submits an emergency message and optional photo",
        "Backend resolves identity from the authenticated account",
        "Role-scoped retrieval shows only what each authority should see",
        "Safety intelligence adds context, severity, and escalation targeting",
        "Incident history becomes searchable and auditable",
    ]:
        p = tf.add_paragraph()
        p.text = item
        style_paragraph(p, font_size=13.2, color=BLACK, bullet=True)

    if MAP_IMAGE.exists():
        slide.shapes.add_picture(str(MAP_IMAGE), Inches(6.35), Inches(1.25), width=Inches(6.3), height=Inches(3.55))
    else:
        fallback = add_rounded_card(slide, Inches(6.35), Inches(1.25), Inches(6.3), Inches(3.55), NAVY)
        tf = fallback.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "Campus map asset not found"
        r.font.name = "Aptos Display"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = WHITE

    add_bullet_box(slide, Inches(6.35), Inches(4.95), Inches(6.3), Inches(1.55), "Live location monitoring", [
        "Coordinates are updated in the backend and displayed on a Leaflet-based map.",
        "Role filters keep admin, HOD, coordinator, and teacher views scoped appropriately.",
    ], title_color=NAVY, body_color=BLACK, fill=WHITE)
    add_footer(slide, "SmartCampus conference deck | safety and location intelligence")


def security_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_header(slide, "Security and Production Hardening", "The backend has a clear security posture, not just feature endpoints")

    controls = [
        "Prepared statements via DB helper",
        "Input validation and sanitization",
        "Bearer-token authentication",
        "CORS, security headers, and rate limiting",
        "bcrypt password hashing and generic errors",
        "Web-server hardening and restricted file access",
    ]
    panel = add_rounded_card(slide, Inches(0.7), Inches(1.35), Inches(5.65), Inches(5.15), NAVY)
    tf = panel.text_frame
    tf.clear()
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Implemented controls"
    r.font.name = "Aptos Display"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WHITE
    for item in controls:
        p = tf.add_paragraph()
        p.text = item
        style_paragraph(p, font_size=13.2, color=RGBColor(238, 242, 247), bullet=True)

    add_bullet_box(slide, Inches(6.65), Inches(1.35), Inches(6.0), Inches(2.25), "Production-readiness signals", [
        "Documentation covers environment setup, deployment, and security checks.",
        "Netlify deployment is defined for the frontend with a separate PHP backend.",
        "The repo includes a pre-production security checklist and operational guidance.",
    ], title_color=NAVY, body_color=BLACK, fill=MIST)

    add_bullet_box(slide, Inches(6.65), Inches(3.85), Inches(6.0), Inches(2.65), "Notable follow-up areas", [
        "Broader prepared-statement adoption is still a consistency task across all endpoints.",
        "Some frontend pages still have lint issues to clean up before stricter release gates.",
        "Bundle size suggests further code-splitting opportunity.",
    ], title_color=NAVY, body_color=BLACK, fill=SAND)

    add_footer(slide, "SmartCampus conference deck | security posture")


def deployment_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MIST)
    add_header(slide, "Deployment and Maturity", "The system already has a credible path to production and demonstration")

    left = add_rounded_card(slide, Inches(0.65), Inches(1.3), Inches(6.1), Inches(4.95), WHITE)
    tf = left.text_frame
    tf.clear()
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Deployment model"
    r.font.name = "Aptos Display"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = NAVY
    for item in [
        "Frontend: Netlify static hosting from the Vite build",
        "Backend: PHP app on shared hosting, VPS, or managed PHP runtime",
        "Database: MySQL schema named smartcampus",
        "Auth and API traffic: HTTPS + bearer tokens + environment variables",
        "The repository includes deployment notes and integration guides",
    ]:
        p = tf.add_paragraph()
        p.text = item
        style_paragraph(p, font_size=13.2, color=BLACK, bullet=True)

    right = add_rounded_card(slide, Inches(6.95), Inches(1.3), Inches(5.75), Inches(4.95), NAVY)
    tf = right.text_frame
    tf.clear()
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Why this is conference-ready"
    r.font.name = "Aptos Display"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WHITE
    for item in [
        "Clear narrative: one platform, multiple roles, one operational view",
        "Demonstrable deep work: security, attendance integrity, and safety workflows",
        "Practical deployment story: static frontend + modular PHP backend",
        "Extensible architecture for future services and policy automation",
    ]:
        p = tf.add_paragraph()
        p.text = item
        style_paragraph(p, font_size=13.2, color=RGBColor(236, 240, 247), bullet=True)

    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(6.35), Inches(9.3), Inches(0.45))
    footer.fill.solid()
    footer.fill.fore_color.rgb = SAND
    footer.line.color.rgb = SAND
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "The project combines implemented workflows, deployment notes, and security guidance into a coherent system story."
    r.font.name = "Aptos"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = NAVY

    add_footer(slide, "SmartCampus conference deck | deployment and maturity")


def close_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NAVY)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MIDNIGHT
    bar.line.color.rgb = MIDNIGHT
    title = slide.shapes.add_textbox(Inches(0.65), Inches(1.0), Inches(6.5), Inches(0.8))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Closing message"
    r.font.name = "Aptos Display"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WHITE
    body = slide.shapes.add_textbox(Inches(0.65), Inches(1.7), Inches(7.0), Inches(2.3))
    tf = body.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "SmartCampus demonstrates how a university can unify academic, operational, and safety workflows into a single role-aware platform."
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "The strongest conference story is not just the breadth of modules, but the integrity model behind them: token-based access, scoped data visibility, layered attendance verification, and safety workflows that can be audited."
    style_paragraph(p2, font_size=15, color=RGBColor(223, 231, 242))

    panel = add_rounded_card(slide, Inches(8.0), Inches(1.45), Inches(4.35), Inches(3.65), WHITE)
    tf = panel.text_frame
    tf.clear()
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Suggested next steps"
    r.font.name = "Aptos Display"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = NAVY
    for item in [
        "Trim a few frontend lint issues",
        "Split large bundles for better performance",
        "Standardize endpoint validation patterns",
        "Add a short live demo or screenshot appendix",
    ]:
        p = tf.add_paragraph()
        p.text = item
        style_paragraph(p, font_size=13, color=BLACK, bullet=True)

    quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(4.55), Inches(11.7), Inches(1.1))
    quote.fill.solid()
    quote.fill.fore_color.rgb = MIDNIGHT
    quote.line.color.rgb = GOLD
    tf = quote.text_frame
    tf.clear()
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "A unified campus platform is most valuable when it makes everyday operations visible, verifiable, and role-aware."
    r.font.name = "Aptos Display"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE

    add_tag(slide, "Thank you", Inches(5.55), Inches(6.1), 1.8, color=GOLD, text_color=NAVY)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "SmartCampus Conference Deck"
    prs.core_properties.subject = "Conference presentation generated from SmartCampus project analysis"
    prs.core_properties.author = "GitHub Copilot"

    cover_slide(prs)
    problem_slide(prs)
    architecture_slide(prs)
    modules_slide(prs)
    attendance_slide(prs)
    safety_slide(prs)
    security_slide(prs)
    deployment_slide(prs)
    close_slide(prs)

    prs.save(str(OUTPUT))


if __name__ == "__main__":
    build_presentation()